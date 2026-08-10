#!/usr/bin/env python3
"""Extract local document text into the Session Workspace."""

from __future__ import annotations

import argparse
import json
import subprocess
from pathlib import Path
from typing import Any

from cxba_material_common import (
    MaterialToolError,
    atomic_write_text,
    load_material,
    sanitized_name,
    write_json,
)


def decode_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return raw.decode(encoding)
        except UnicodeError:
            pass
    raise MaterialToolError("Text encoding is not UTF-8 or GB18030")


def extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        parts.append(f"\n\n--- page {index} ---\n{text}")
        pages.append({"page": index, "characters": len(text)})
    return "".join(parts).lstrip(), {"pages": pages}


def extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    from docx import Document

    document = Document(path)
    blocks = [paragraph.text for paragraph in document.paragraphs]
    table_count = 0
    for table_index, table in enumerate(document.tables, start=1):
        table_count += 1
        blocks.append(f"\n--- table {table_index} ---")
        blocks.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
    return "\n".join(blocks), {"paragraphCount": len(document.paragraphs), "tableCount": table_count}


def extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(path)
    blocks = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        blocks.append(f"\n--- slide {slide_index} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                blocks.append(shape.text)
    return "\n".join(blocks).lstrip(), {"slideCount": len(presentation.slides)}


def extract_rtf(path: Path) -> tuple[str, dict[str, Any]]:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(decode_text(path)), {}


def extract_doc(path: Path) -> tuple[str, dict[str, Any]]:
    errors = []
    for command in ("antiword", "catdoc"):
        result = subprocess.run([command, str(path)], capture_output=True, check=False)
        if result.returncode == 0 and result.stdout.strip():
            for encoding in ("utf-8", "gb18030"):
                try:
                    return result.stdout.decode(encoding), {"extractor": command}
                except UnicodeError:
                    pass
        errors.append(f"{command}: exit={result.returncode}")
    raise MaterialToolError("Legacy Word extraction failed: " + ", ".join(errors))


def extract(path: Path) -> tuple[str, dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".rtf":
        return extract_rtf(path)
    if suffix == ".doc":
        return extract_doc(path)
    if suffix in {".txt", ".md", ".json", ".csv", ".tsv", ".xml", ".html"}:
        return decode_text(path), {}
    raise MaterialToolError(f"Unsupported document format: {suffix or '<none>'}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--catalog", required=True, type=Path)
    parser.add_argument("--material-id", required=True)
    parser.add_argument("--output-dir", required=True, type=Path)
    args = parser.parse_args()

    entry, path = load_material(args.catalog, args.material_id)
    text, details = extract(path)
    args.output_dir.mkdir(parents=True, exist_ok=True)
    stem = sanitized_name(args.material_id)
    text_path = args.output_dir / f"{stem}.txt"
    metadata_path = args.output_dir / f"{stem}.json"
    atomic_write_text(text_path, text)
    write_json(
        metadata_path,
        {
            "materialId": args.material_id,
            "relativePath": entry["relativePath"],
            "textPath": str(text_path),
            "characters": len(text),
            **details,
        },
    )
    print(json.dumps({"status": "completed", "textPath": str(text_path), "metadataPath": str(metadata_path)}))


if __name__ == "__main__":
    try:
        main()
    except MaterialToolError as exc:
        raise SystemExit(f"document_extract_failed: {exc}") from exc
